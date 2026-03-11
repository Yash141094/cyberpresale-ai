"""
PowerPoint export for CyberPresales AI
Modern proposal format — structured sections, tables, architecture diagram, visual elements
"""
import os
import re
import tempfile
from datetime import datetime


def _clean_md(text, max_chars=700):
    if not text:
        return "Analysis not yet generated."
    lines = []
    for line in text.split('\n'):
        line = line.strip()
        if not line:
            continue
        line = re.sub(r'^#{1,4}\s*', '', line)
        line = re.sub(r'\*\*(.*?)\*\*', r'\1', line)
        line = re.sub(r'\*(.*?)\*', r'\1', line)
        line = re.sub(r'^\|.*\|$', '', line)
        line = re.sub(r'^[-*]\s+', '* ', line)
        if line:
            lines.append(line)
    result = '\n'.join(lines)
    return result[:max_chars] + ('...' if len(result) > max_chars else '')


def _parse_sections(text, max_sections=6):
    if not text:
        return []
    sections = []
    parts = re.split(r'^##\s+', text, flags=re.MULTILINE)
    for part in parts[1:max_sections+1]:
        lines = part.strip().split('\n')
        heading = lines[0].strip()
        body = _clean_md('\n'.join(lines[1:]), max_chars=400)
        sections.append((heading, body))
    return sections


def generate_ppt(session_state):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        import subprocess
        subprocess.run(["pip", "install", "python-pptx", "--break-system-packages", "-q"])
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

    NAVY   = RGBColor(0x0d, 0x1b, 0x2a)
    DARK   = RGBColor(0x11, 0x24, 0x38)
    PANEL  = RGBColor(0x1a, 0x2f, 0x45)
    ACCENT = RGBColor(0x00, 0x8b, 0xd8)
    GREEN  = RGBColor(0x00, 0xc9, 0x96)
    AMBER  = RGBColor(0xf5, 0x9e, 0x0b)
    RED    = RGBColor(0xef, 0x44, 0x44)
    PINK   = RGBColor(0xec, 0x48, 0x99)
    PURPLE = RGBColor(0x8b, 0x5c, 0xf6)
    WHITE  = RGBColor(0xff, 0xff, 0xff)
    LIGHT  = RGBColor(0xc8, 0xd8, 0xe8)
    MUTED  = RGBColor(0x6b, 0x88, 0xa0)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def get(key):
        if hasattr(session_state, key):
            return getattr(session_state, key)
        if hasattr(session_state, '__getitem__'):
            return session_state.get(key)
        return None

    def R(slide, x, y, w, h, fill):
        s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
        return s

    def T(slide, x, y, w, h, text, size=13, color=WHITE, bold=False,
          align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf  = box.text_frame; tf.word_wrap = True
        p   = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = str(text)
        run.font.size = Pt(size); run.font.color.rgb = color
        run.font.bold = bold; run.font.italic = italic; run.font.name = font
        return box

    def bg(slide):       R(slide, 0, 0, 13.33, 7.5, DARK)
    def lbar(slide, c):  R(slide, 0, 0, 0.07, 7.5, c)
    def hband(slide):    R(slide, 0, 0, 13.33, 1.15, NAVY)
    def footer(slide):
        T(slide, 0.15, 7.18, 9, 0.28,
          f"CyberPresales AI  |  Confidential  |  {datetime.now().strftime('%B %Y')}",
          size=7, color=MUTED)

    def divider(slide, title, sub="", color=ACCENT):
        bg(slide); R(slide, 0, 0, 0.4, 7.5, color)
        R(slide, 0.4, 3.0, 12.93, 0.04, color)
        T(slide, 0.7, 3.1, 11, 1.1, title, size=40, bold=True, color=WHITE)
        if sub: T(slide, 0.7, 4.3, 10, 0.5, sub, size=15, color=LIGHT, italic=True)
        footer(slide)

    def card(slide, x, y, w, h, title, body, tc=ACCENT):
        R(slide, x, y, w, h, PANEL)
        T(slide, x+0.12, y+0.1, w-0.2, 0.3, title, size=9, color=tc, bold=True)
        T(slide, x+0.12, y+0.48, w-0.2, h-0.58, body, size=8, color=LIGHT)

    def stat(slide, x, y, val, lbl, color=ACCENT):
        R(slide, x, y, 2.5, 1.1, PANEL); R(slide, x, y, 2.5, 0.06, color)
        T(slide, x+0.1, y+0.1, 2.3, 0.55, val, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        T(slide, x+0.1, y+0.72, 2.3, 0.32, lbl, size=9, color=MUTED, align=PP_ALIGN.CENTER)

    # ─── SLIDE 1: COVER ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s)
    R(s, 0, 0, 13.33, 0.12, ACCENT); R(s, 0, 7.38, 13.33, 0.12, ACCENT)
    R(s, 0, 0.12, 4.2, 7.26, NAVY); R(s, 4.2, 0.12, 0.04, 7.26, ACCENT)
    # Dynamic cover type label from detected RFP type
    domains_meta = get('domains') or {}
    rfp_type_cover = domains_meta.get('rfp_type', 'IT SERVICES') if isinstance(domains_meta, dict) else 'IT SERVICES'
    cover_label = rfp_type_cover.upper()
    T(s, 0.3, 0.5, 3.6, 0.45, cover_label, size=10 if len(cover_label) > 15 else 12, color=ACCENT, bold=True)
    T(s, 0.3, 0.95, 3.6, 1.2, "Presales\nIntelligence\nReport", size=28, bold=True, color=WHITE)
    fn = (get('file_name') or 'RFP Document')[:45]
    for lbl, val, yp in [("Document", fn, 2.5), ("Generated", datetime.now().strftime('%d %B %Y'), 3.4),
                          ("Author", "Yash Mehrotra", 4.2), ("Platform", "CyberPresales AI v3", 5.0)]:
        T(s, 0.3, yp, 3.6, 0.28, lbl, size=8, color=MUTED)
        T(s, 0.3, yp+0.3, 3.6, 0.38, val, size=10, color=LIGHT if lbl != "Platform" else ACCENT, bold=(lbl=="Platform"))
    T(s, 4.6, 1.2, 8.4, 1.9, "AI-Powered\nPresales Intelligence", size=38, bold=True, color=WHITE)
    T(s, 4.6, 3.2, 8.0, 0.45, "Transforming RFP analysis into winning proposals", size=14, color=LIGHT, italic=True)
    n = sum(1 for k in ["customer_brief","pain_analysis","solution_rec","competitive","product_map","exec_summary"] if get(k))
    stat(s, 4.6, 4.3, str(n), "Modules Generated", ACCENT)
    stat(s, 7.3, 4.3, "9", "Intelligence Modules", GREEN)
    stat(s, 10.0, 4.3, "v3", "Platform", AMBER)

    # ─── SLIDE 2: CONTENTS ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s); lbar(s, ACCENT); hband(s)
    T(s, 0.3, 0.18, 10, 0.75, "Proposal Contents", size=28, bold=True, color=WHITE)
    T(s, 0.3, 0.88, 10, 0.25, "AI-generated modules included in this report", size=10, color=MUTED, italic=True)
    toc = [("01","Executive Summary","Board-ready proposal narrative",bool(get('exec_summary'))),
           ("02","Customer Intelligence","Who they are · Entry points · Decision makers",bool(get('customer_brief'))),
           ("03","Pain Point Analysis","Surface requirements vs underlying business pains",bool(get('pain_analysis'))),
           ("04","Solution Recommendation","Architecture · Phasing · POC strategy",bool(get('solution_rec'))),
           ("05","Competitive Intelligence","Likely bidders · Counter-strategy · Win conditions",bool(get('competitive'))),
           ("06","Product & Vendor Mapping","Domain-by-domain · Primary & alternative vendors",bool(get('product_map'))),
           ("07","Solution Architecture","Visual architecture layers · Integration model",True)]
    y = 1.3
    for num, nm, desc, done in toc:
        col = GREEN if done else MUTED
        R(s, 0.25, y, 0.55, 0.48, PANEL)
        T(s, 0.25, y+0.06, 0.55, 0.35, num, size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
        T(s, 0.95, y+0.02, 7.5, 0.26, nm, size=12, bold=True, color=WHITE if done else MUTED)
        T(s, 0.95, y+0.27, 7.5, 0.22, desc, size=8, color=MUTED, italic=True)
        T(s, 10.5, y+0.1, 2.6, 0.3, "Included" if done else "Not Generated",
          size=10, color=GREEN if done else MUTED, align=PP_ALIGN.RIGHT)
        y += 0.68
    footer(s)

    # ─── SLIDE 3-4: EXECUTIVE SUMMARY ─────────────────────────────────────────
    s = prs.slides.add_slide(blank); divider(s, "Executive Summary", "Strategic overview for C-level stakeholders", ACCENT)
    if get('exec_summary'):
        s = prs.slides.add_slide(blank); bg(s); lbar(s, ACCENT); hband(s)
        T(s, 0.3, 0.18, 10, 0.75, "Executive Summary", size=26, bold=True, color=WHITE)
        T(s, 0.3, 1.2, 12.7, 6.0, _clean_md(get('exec_summary'), 900), size=11, color=LIGHT)
        footer(s)

    # ─── SLIDE 5-6: CUSTOMER BRIEF ────────────────────────────────────────────
    s = prs.slides.add_slide(blank); divider(s, "Customer Intelligence", "Understanding who we are selling to", GREEN)
    if get('customer_brief'):
        s = prs.slides.add_slide(blank); bg(s); lbar(s, GREEN); hband(s)
        T(s, 0.3, 0.18, 10, 0.75, "Customer Intelligence Brief", size=26, bold=True, color=WHITE)
        T(s, 0.3, 1.2, 12.7, 6.0, _clean_md(get('customer_brief'), 900), size=11, color=LIGHT)
        footer(s)

    # ─── SLIDE 7-8: PAIN ANALYSIS ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank); divider(s, "Pain Point Analysis", "Surface requirements vs underlying business pains", AMBER)
    if get('pain_analysis'):
        s = prs.slides.add_slide(blank); bg(s); lbar(s, AMBER); hband(s)
        T(s, 0.3, 0.18, 10, 0.75, "Pain Point & Requirements Analysis", size=26, bold=True, color=WHITE)
        secs = _parse_sections(get('pain_analysis'), 4)
        if secs:
            pos = [(0.25,1.2,6.3,2.8),(6.8,1.2,6.3,2.8),(0.25,4.2,6.3,2.8),(6.8,4.2,6.3,2.8)]
            cols = [AMBER, RED, ACCENT, GREEN]
            for i,(h,b) in enumerate(secs[:4]): card(s, *pos[i], h, b, cols[i])
        else:
            T(s, 0.3, 1.2, 12.7, 6.0, _clean_md(get('pain_analysis'), 900), size=11, color=LIGHT)
        footer(s)

    # ─── SLIDE 9-10: SOLUTION REC ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank); divider(s, "Solution Recommendation", "What to propose · Architecture · POC strategy", ACCENT)
    if get('solution_rec'):
        s = prs.slides.add_slide(blank); bg(s); lbar(s, ACCENT); hband(s)
        T(s, 0.3, 0.18, 10, 0.75, "Solution Recommendation", size=26, bold=True, color=WHITE)
        secs = _parse_sections(get('solution_rec'), 3)
        if secs:
            pos = [(0.25,1.2,12.83,1.9),(0.25,3.25,6.3,3.8),(6.8,3.25,6.3,3.8)]
            cols = [ACCENT, GREEN, AMBER]
            for i,(h,b) in enumerate(secs[:3]): card(s, *pos[i], h, b, cols[i])
        else:
            T(s, 0.3, 1.2, 12.7, 6.0, _clean_md(get('solution_rec'), 900), size=11, color=LIGHT)
        footer(s)

    # ─── SLIDE 11: ARCHITECTURE DIAGRAM ───────────────────────────────────────
    # Detect RFP type first — used for both subtitle and layer selection
    domains_data      = get('domains') or {}
    rfp_type_detected = (domains_data.get('rfp_type','') if isinstance(domains_data,dict) else '') or 'Cybersecurity'

    # Integration bar text adapts to domain
    integration_bar_map = {
        'Cybersecurity':             "UNIFIED SECURITY DATA LAKE  |  SOAR ORCHESTRATION  |  SINGLE PANE OF GLASS",
        'Infrastructure Services':   "UNIFIED MONITORING PLATFORM  |  ITSM INTEGRATION  |  SINGLE PANE OF GLASS",
        'Application Managed Services': "ITSM PLATFORM (SERVICENOW / JIRA)  |  CI/CD PIPELINE  |  UNIFIED REPORTING",
        'End User Computing':        "ITSM PLATFORM  |  DEX ANALYTICS  |  SELF-SERVICE PORTAL  |  KNOWLEDGE BASE",
        'Digital Workplace':         "MICROSOFT 365 FABRIC  |  PURVIEW GOVERNANCE  |  UNIFIED ADMIN CENTRE",
        'Data & Analytics':          "UNIFIED DATA PLATFORM  |  DATA GOVERNANCE LAYER  |  MLOPS ORCHESTRATION",
        'Multi-Tower':               "INTEGRATED ITSM (SERVICENOW)  |  UNIFIED REPORTING  |  GOVERNANCE LAYER",
    }
    integration_bar = integration_bar_map.get(rfp_type_detected, integration_bar_map['Multi-Tower'])

    s = prs.slides.add_slide(blank); bg(s); lbar(s, ACCENT); hband(s)
    T(s, 0.3, 0.18, 10, 0.75, "Solution Architecture", size=26, bold=True, color=WHITE)
    T(s, 0.3, 0.85, 12, 0.28, f"{rfp_type_detected} — integrated delivery architecture", size=10, color=MUTED, italic=True)
    R(s, 0.25, 1.2, 12.83, 0.38, NAVY)
    T(s, 0.3, 1.25, 12.7, 0.3, integration_bar, size=8, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Architecture templates per RFP type
    arch_templates = {
        'Cybersecurity': [
            ("USERS & DEVICES",      "EDR / XDR  ·  MDM  ·  Endpoint Protection  ·  Threat Hunting  ·  DLP",         AMBER),
            ("IDENTITY & ZERO TRUST","IAM  ·  PAM  ·  MFA  ·  Conditional Access  ·  ZTNA  ·  SSO",                  ACCENT),
            ("NETWORK SECURITY",     "NGFW  ·  IDS/IPS  ·  WAF  ·  Microsegmentation  ·  SD-WAN  ·  DDoS",           PURPLE),
            ("CLOUD SECURITY",       "CSPM  ·  CWPP  ·  CNAPP  ·  Container Security  ·  API Protection",             GREEN),
            ("SECURITY OPERATIONS",  "SIEM  ·  SOAR  ·  UEBA  ·  Threat Intel  ·  24x7 SOC  ·  GRC",                 RED),
        ],
        'Infrastructure Services': [
            ("COMPUTE & VIRTUALISATION","Servers  ·  VMware  ·  Hyper-V  ·  Container Platform  ·  HCI",              AMBER),
            ("NETWORK & CONNECTIVITY",  "LAN/WAN  ·  SD-WAN  ·  Load Balancers  ·  Firewall  ·  DNS/DHCP",            ACCENT),
            ("STORAGE & BACKUP",        "SAN/NAS  ·  Backup & Recovery  ·  DR  ·  Tape/Cloud Archive",                PURPLE),
            ("CLOUD PLATFORM",          "AWS / Azure / GCP  ·  Cloud Migration  ·  Cost Optimisation  ·  IaC",        GREEN),
            ("MONITORING & OPERATIONS", "NOC  ·  ITSM  ·  Observability  ·  Automation  ·  Capacity Planning",        RED),
        ],
        'Application Managed Services': [
            ("L1 / FUNCTIONAL SUPPORT", "User Queries  ·  Access Management  ·  How-To  ·  Password Reset",           AMBER),
            ("L2 / TECHNICAL SUPPORT",  "Bug Triage  ·  Config Changes  ·  Performance Issues  ·  Patch Application", ACCENT),
            ("L3 / ENGINEERING",        "Root Cause Analysis  ·  Code Fixes  ·  Architecture Review  ·  Hot Fixes",   PURPLE),
            ("RELEASE & DEVOPS",        "CI/CD Pipeline  ·  Release Management  ·  Test Automation  ·  Deployment",   GREEN),
            ("GOVERNANCE & REPORTING",  "SLA Reporting  ·  ITSM (ServiceNow/JIRA)  ·  Audit  ·  Documentation",       RED),
        ],
        'End User Computing': [
            ("SERVICE DESK",          "L1 Helpdesk  ·  Incident Logging  ·  Password Reset  ·  24x7 Coverage",        AMBER),
            ("DESKTOP & DEVICE MGMT", "Imaging  ·  Patch Management  ·  MDM  ·  Asset Lifecycle  ·  SCCM/Intune",     ACCENT),
            ("PRODUCTIVITY PLATFORM", "M365 Administration  ·  Teams  ·  SharePoint  ·  OneDrive  ·  Exchange",       PURPLE),
            ("VDI & REMOTE ACCESS",   "Citrix / AVD  ·  VPN  ·  Remote Support  ·  BYOD Policy",                      GREEN),
            ("ITSM & REPORTING",      "ServiceNow  ·  SLA Dashboards  ·  CSAT  ·  Knowledge Base  ·  Automation",     RED),
        ],
        'Multi-Tower': [
            ("SERVICE DESK & EUC",    "Helpdesk L1/L2  ·  Desktop  ·  Device Management  ·  M365",                    AMBER),
            ("INFRASTRUCTURE",        "DC Operations  ·  Cloud  ·  Network  ·  Storage  ·  Backup",                   ACCENT),
            ("APPLICATION SERVICES",  "AMS L2/L3  ·  Release Management  ·  DevOps  ·  Testing",                      PURPLE),
            ("CYBERSECURITY",         "SOC  ·  EDR  ·  Identity  ·  Cloud Security  ·  GRC",                          GREEN),
            ("GOVERNANCE & ITSM",     "ServiceNow  ·  SLA Reporting  ·  CSAT  ·  Continuous Improvement  ·  Audit",   RED),
        ],
    }

    layers = arch_templates.get(rfp_type_detected) or arch_templates['Multi-Tower']
    y = 1.65
    for nm, tools, col in layers:
        R(s, 0.25, y, 12.83, 0.95, PANEL); R(s, 0.25, y, 0.09, 0.95, col)
        T(s, 0.42, y+0.08, 2.8, 0.35, nm, size=9, bold=True, color=col)
        T(s, 0.42, y+0.5,  12.2, 0.38, tools, size=9, color=LIGHT)
        y += 1.0
    footer(s)

    # ─── SLIDE 12-13: COMPETITIVE ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank); divider(s, "Competitive Intelligence", "Who is bidding · Where we win · Counter-strategy", PINK)
    if get('competitive'):
        s = prs.slides.add_slide(blank); bg(s); lbar(s, PINK); hband(s)
        T(s, 0.3, 0.18, 10, 0.75, "Competitive Intelligence", size=26, bold=True, color=WHITE)
        secs = _parse_sections(get('competitive'), 4)
        if secs:
            pos = [(0.25,1.2,6.3,2.8),(6.8,1.2,6.3,2.8),(0.25,4.2,6.3,2.8),(6.8,4.2,6.3,2.8)]
            cols = [PINK, RED, ACCENT, GREEN]
            for i,(h,b) in enumerate(secs[:4]): card(s, *pos[i], h, b, cols[i])
        else:
            T(s, 0.3, 1.2, 12.7, 6.0, _clean_md(get('competitive'), 900), size=11, color=LIGHT)
        footer(s)

    # ─── SLIDE 14-15: PRODUCT MAPPING ─────────────────────────────────────────
    s = prs.slides.add_slide(blank); divider(s, "Product & Vendor Mapping", "Domain-by-domain recommendations", PURPLE)
    if get('product_map'):
        s = prs.slides.add_slide(blank); bg(s); lbar(s, PURPLE); hband(s)
        T(s, 0.3, 0.18, 10, 0.75, "Product & Vendor Recommendations", size=26, bold=True, color=WHITE)
        secs = _parse_sections(get('product_map'), 6)
        if secs:
            y = 1.18
            rh = 0.87
            for i,(domain, content) in enumerate(secs[:6]):
                bg_c = PANEL if i%2==0 else NAVY
                R(s, 0.25, y, 12.83, rh, bg_c); R(s, 0.25, y, 2.9, rh, DARK)
                T(s, 0.35, y+0.15, 2.6, 0.6, domain, size=9, bold=True, color=PURPLE)
                T(s, 3.25, y+0.08, 9.6, rh-0.12,
                  content[:200]+('...' if len(content)>200 else ''), size=8, color=LIGHT)
                y += rh + 0.04
        else:
            T(s, 0.3, 1.2, 12.7, 6.0, _clean_md(get('product_map'), 900), size=11, color=LIGHT)
        footer(s)

    # ─── SLIDE 16: WIN STRATEGY ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s); lbar(s, GREEN); hband(s)
    T(s, 0.3, 0.18, 10, 0.75, "Deal Win Strategy", size=26, bold=True, color=WHITE)
    T(s, 0.3, 0.85, 12, 0.28, "Key elements required to win this engagement", size=10, color=MUTED, italic=True)
    win = [
        ("PROPOSAL MUST-HAVES",
         "* Reference customers in same industry\n* Compliance mapped to each requirement\n* TCO analysis vs point solutions\n* Implementation timeline with milestones\n* Named resources and local team",
         GREEN),
        ("POC STRATEGY",
         "* Focus on the highest-pain requirement\n* Demonstrate integration with existing tools\n* Show time-to-detect improvement vs current\n* Deliver a live dashboard the customer keeps\n* Present insights during POC, not after",
         ACCENT),
        ("COMMERCIAL STRATEGY",
         "* Lead with platform value not product cost\n* Anchor against breach cost and penalties\n* Offer phased commercials to reduce commitment\n* Bundle SLAs into base price\n* Quantify 3-year TCO comparison",
         AMBER),
    ]
    x = 0.25
    for title, body, col in win:
        R(s, x, 1.3, 4.22, 5.9, PANEL); R(s, x, 1.3, 4.22, 0.06, col)
        T(s, x+0.15, 1.45, 3.9, 0.38, title, size=9, bold=True, color=col)
        T(s, x+0.15, 1.9,  3.9, 5.1,  body,  size=10, color=LIGHT)
        x += 4.44
    footer(s)

    # ─── SLIDE 17: DOMAIN SCOPE PIE CHART ────────────────────────────────────
    try:
        import io, math
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import matplotlib.patches as mpatches

        rfp_lower_chart = (get('rfp_text') or '').lower()
        tower_kws = {
            'Cybersecurity':             ['siem','soc','edr','threat','grc','firewall','ztna','dlp','soar'],
            'Infrastructure & DC Ops':   ['server','data centre','dc','compute','vmware','hci','storage','san','nas','backup','dr'],
            'App Managed Services':      ['sap','erp','ams','l1','l2','l3','application support','release','devops','jira'],
            'End User Computing':        ['helpdesk','desktop','endpoint','euc','m365','intune','vdi','citrix','service desk','patch'],
            'Digital Workplace':         ['sharepoint','teams','intranet','collaboration','purview','onedrive','power platform'],
            'Cloud & Migration':         ['cloud','aws','azure','gcp','migration','lift-and-shift','kubernetes','container','finops'],
            'Data & Analytics':          ['data lake','analytics','bi','etl','data warehouse','mlops','power bi','tableau','databricks'],
            'Networking':                ['wan','lan','bgp','routing','switching','mpls','sd-wan','noc','wireless'],
        }
        w = {t: sum(rfp_lower_chart.count(kw) for kw in kws) for t, kws in tower_kws.items()}
        w = {t: v for t, v in w.items() if v > 0}
        if not w:
            domains_from_classify = get('domains') or {}
            for d_ in (domains_from_classify.get('detected_domains', []) if isinstance(domains_from_classify, dict) else []):
                w[d_] = 10
        if not w:
            w = {'IT Services': 1}

        chart_labels = list(w.keys())
        chart_values = list(w.values())
        total_w = sum(chart_values)
        chart_pcts = [v/total_w*100 for v in chart_values]

        pie_colors = ['#c2550a','#ea7c2b','#d97706','#166534','#1e40af','#6b21a8','#0e7490','#be185d'][:len(chart_labels)]

        fig_pie, ax_pie = plt.subplots(figsize=(8, 4.5), facecolor='#0d1b2a')
        ax_pie.set_facecolor('#0d1b2a')
        wedges, texts, autotexts = ax_pie.pie(
            chart_values, labels=None, autopct='%1.0f%%',
            colors=pie_colors, startangle=90,
            wedgeprops=dict(width=0.55, edgecolor='#0d1b2a', linewidth=2),
            pctdistance=0.75,
        )
        for at in autotexts:
            at.set_fontsize(9); at.set_color('white'); at.set_fontweight('bold')
        legend_patches = [mpatches.Patch(color=pie_colors[i], label=f"{chart_labels[i]} ({chart_pcts[i]:.0f}%)")
                          for i in range(len(chart_labels))]
        ax_pie.legend(handles=legend_patches, loc='center left', bbox_to_anchor=(1.0, 0.5),
                      fontsize=8.5, frameon=False, labelcolor='#c8d8e8')
        ax_pie.set_title('Scope Distribution by Service Tower', color='white', fontsize=12, pad=12, fontweight='bold')
        plt.tight_layout()
        buf_pie = io.BytesIO(); plt.savefig(buf_pie, format='png', dpi=150, bbox_inches='tight', facecolor='#0d1b2a'); buf_pie.seek(0); plt.close()

        s = prs.slides.add_slide(blank); bg(s); lbar(s, ACCENT); hband(s)
        T(s, 0.3, 0.18, 10, 0.75, "Domain Scope Distribution", size=26, bold=True, color=WHITE)
        T(s, 0.3, 0.85, 12, 0.28, "Keyword-weighted breakdown of RFP coverage across IT service towers", size=10, color=MUTED, italic=True)
        from pptx.util import Inches as _In
        pic = s.shapes.add_picture(buf_pie, _In(0.3), _In(1.2), _In(12.5), _In(5.9))
        footer(s)
    except Exception as _chart_err:
        pass  # chart slide optional — skip silently if matplotlib unavailable

    # ─── SLIDE 18: COMPETITIVE BAR CHART (if competitive data present) ────────
    try:
        comp_text = get('competitive') or ''
        if comp_text:
            import re as _re
            # Extract competitor names and threat levels from table in competitive text
            comp_rows = _re.findall(r'\|\s*([A-Za-z &]+)\s*\|\s*(High|Medium|Med|Low)\s*\|', comp_text, _re.IGNORECASE)
            if len(comp_rows) >= 2:
                comp_names  = [r[0].strip() for r in comp_rows[:6]]
                threat_map  = {'high': 3, 'medium': 2, 'med': 2, 'low': 1}
                threat_vals = [threat_map.get(r[1].strip().lower(), 2) for r in comp_rows[:6]]
                bar_colors  = ['#ef4444' if v==3 else '#f59e0b' if v==2 else '#34d399' for v in threat_vals]

                fig_bar, ax_bar = plt.subplots(figsize=(9, 3.8), facecolor='#0d1b2a')
                ax_bar.set_facecolor('#0d1b2a')
                bars = ax_bar.barh(comp_names, threat_vals, color=bar_colors, height=0.5, edgecolor='#0d1b2a')
                ax_bar.set_xlim(0, 3.5)
                ax_bar.set_xticks([1, 2, 3]); ax_bar.set_xticklabels(['Low', 'Medium', 'High'], color='#9aa0b4', fontsize=9)
                ax_bar.tick_params(axis='y', colors='#c8d8e8', labelsize=10)
                ax_bar.spines['top'].set_visible(False); ax_bar.spines['right'].set_visible(False)
                ax_bar.spines['bottom'].set_color('#252b38'); ax_bar.spines['left'].set_color('#252b38')
                ax_bar.set_title('Competitive Threat Assessment', color='white', fontsize=12, pad=10, fontweight='bold')
                ax_bar.axvline(x=2.5, color='#ef4444', linestyle='--', alpha=0.4, linewidth=1)
                plt.tight_layout()
                buf_bar = io.BytesIO(); plt.savefig(buf_bar, format='png', dpi=150, bbox_inches='tight', facecolor='#0d1b2a'); buf_bar.seek(0); plt.close()

                s = prs.slides.add_slide(blank); bg(s); lbar(s, PINK); hband(s)
                T(s, 0.3, 0.18, 10, 0.75, "Competitive Threat Assessment", size=26, bold=True, color=WHITE)
                T(s, 0.3, 0.85, 12, 0.28, "Relative threat level by competitor for this specific deal", size=10, color=MUTED, italic=True)
                s.shapes.add_picture(buf_bar, _In(0.5), _In(1.2), _In(9.5), _In(5.5))

                # Threat legend
                for i, (name, val) in enumerate(zip(comp_names, threat_vals)):
                    lbl = {3:'HIGH',2:'MEDIUM',1:'LOW'}[val]
                    col = {3:RED,2:AMBER,1:GREEN}[val]
                    T(s, 10.2, 1.5+i*0.9, 2.9, 0.38, f"{name}: {lbl}", size=9, color=col, bold=(val==3))
                footer(s)
    except Exception:
        pass

    # ─── SLIDE 19: IMPLEMENTATION PHASE TIMELINE ─────────────────────────────
    try:
        sol_text = get('solution_rec') or ''
        if sol_text:
            phases = _re.findall(r'Phase\s+([123])[^\n]*\n([^\n#]{20,200})', sol_text)
            if len(phases) >= 2:
                s = prs.slides.add_slide(blank); bg(s); lbar(s, GREEN); hband(s)
                T(s, 0.3, 0.18, 10, 0.75, "Implementation Phasing", size=26, bold=True, color=WHITE)
                T(s, 0.3, 0.85, 12, 0.28, "Proposed delivery phases and key milestones", size=10, color=MUTED, italic=True)

                phase_colors = [GREEN, ACCENT, AMBER]
                ph_labels = [f"Phase {p[0]}" for p in phases[:3]]
                ph_bodies = [p[1][:180] for p in phases[:3]]

                # Timeline bar
                R(s, 0.25, 1.3, 12.83, 0.06, MUTED)
                n = len(ph_labels)
                w_each = 12.83 / n
                for i, (lbl, body, col) in enumerate(zip(ph_labels, ph_bodies, phase_colors)):
                    x = 0.25 + i * w_each
                    R(s, x + 0.05, 1.15, 0.25, 0.35, col)  # milestone dot (rect)
                    T(s, x + 0.05, 1.55, w_each - 0.15, 0.35, lbl, size=11, bold=True, color=col)
                    R(s, x + 0.05, 2.0, w_each - 0.15, 4.8, PANEL)
                    R(s, x + 0.05, 2.0, w_each - 0.15, 0.06, col)
                    T(s, x + 0.18, 2.15, w_each - 0.4, 4.4, body, size=9, color=LIGHT)
                footer(s)
    except Exception:
        pass

    # ─── SLIDE 20: CLOSING ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); bg(s)
    R(s, 0, 0, 13.33, 0.12, ACCENT); R(s, 0, 7.38, 13.33, 0.12, ACCENT)
    R(s, 0, 0.12, 4.2, 7.26, NAVY); R(s, 4.2, 0.12, 0.04, 7.26, ACCENT)
    T(s, 0.3, 1.5,  3.6, 0.35, "POWERED BY",           size=8,  color=MUTED)
    T(s, 0.3, 1.85, 3.6, 0.55, "CyberPresales AI",     size=16, bold=True, color=ACCENT)
    T(s, 0.3, 2.45, 3.6, 0.35, "Groq LLaMA 3.3 70B",  size=10, color=LIGHT)
    T(s, 0.3, 2.85, 3.6, 0.35, "Built by Yash Mehrotra", size=10, color=LIGHT)
    T(s, 0.3, 3.25, 3.6, 0.35, f"v3.0  ·  {datetime.now().strftime('%Y')}", size=10, color=MUTED)
    T(s, 4.6, 2.0,  8.4, 1.4,  "Ready to\nWin This Deal?", size=42, bold=True, color=WHITE)
    T(s, 4.6, 3.6,  8.0, 0.45, "This analysis was generated by CyberPresales AI", size=14, color=LIGHT, italic=True)
    T(s, 4.6, 4.2,  8.0, 0.38, "Every insight is grounded in the RFP you uploaded.", size=11, color=MUTED)
    x = 4.6
    for val, lbl, col in [("9","Modules",ACCENT),("AI","Signal Extract",GREEN),("v3","Platform",AMBER)]:
        stat(s, x, 5.2, val, lbl, col); x += 2.8

    out = os.path.join(tempfile.gettempdir(), "CyberPresales_Proposal.pptx")
    prs.save(out)
    return out
